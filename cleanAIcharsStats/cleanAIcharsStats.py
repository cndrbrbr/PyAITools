#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# pip install python-docx python-pptx openpyxl pymupdf
# py -m pip install python-docx python-pptx openpyxl pymupdf

from __future__ import annotations
import argparse
import sys
import shutil
from pathlib import Path
from typing import Iterable, Tuple, Optional, List
import threading
import queue
from collections import Counter
import unicodedata

# ---------- optionale Libraries ----------
# DOCX
try:
    from docx import Document  # type: ignore
    HAVE_PYTHON_DOCX = True
except Exception:
    HAVE_PYTHON_DOCX = False

# PPTX
try:
    from pptx import Presentation  # type: ignore
    HAVE_PYTHON_PPTX = True
except Exception:
    HAVE_PYTHON_PPTX = False

# XLSX
try:
    import openpyxl  # type: ignore
    HAVE_OPENPYXL = True
except Exception:
    HAVE_OPENPYXL = False

# PDF (Text extrahieren via PyMuPDF)
try:
    import fitz  # type: ignore  # PyMuPDF
    HAVE_PYMUPDF = True
except Exception:
    HAVE_PYMUPDF = False

# ---------- Ersetzungen ----------
REPLACEMENTS = {
    "\u202F": " ",   # Narrow No-Break Space → normales Leerzeichen
    "\u200B": "",    # Zero Width Space → löschen
    "\u2060": "",    # Word Joiner → löschen
    "\u00A0": " ",   # No-Break Space → normales Leerzeichen
    "\u200C": "",    # Zero Width Non-Joiner → löschen
    "\u200D": "",    # Zero Width Joiner → löschen (Achtung bei Emojis)
    "\u2061": "",    # Function Application → löschen
    "\u2062": " ",   # Invisible Times → Leerraum
    "\u2063": " ",   # Invisible Separator → Leerraum
    "\u20AD": " ",   # Kip-Symbol
    "\u2064": " ",   # Invisible Plus → Leerraum
    "\uFEFF": "",    # BOM → löschen
    "\u200E": "",    # LRM
    "\u200F": "",    # RLM
    "\u2028": "\n",  # Line Separator → Zeilenumbruch
    "\u2029": "\n",  # Paragraph Separator → Zeilenumbruch
}

def clean_text(text: str, stats: Optional[Counter[str]] = None) -> str:
    """
    Ersetzt definierte Unicode-Steuer-/Sonderzeichen.
    Wenn stats übergeben wird, werden die Ersetzungen pro Zeichen gezählt.
    """
    for bad, good in REPLACEMENTS.items():
        if bad in text:
            if stats is not None:
                stats[bad] += text.count(bad)
            text = text.replace(bad, good)
    return text

def _format_stats_one_liner(stats: Counter[str], max_items: int = 4) -> str:
    """Kompakt: "Repl: 12 (U+200B:10, U+00A0:2, ...)" """
    total = sum(stats.values())
    if total <= 0:
        return "Repl: 0"
    parts = []
    for ch, cnt in sorted(stats.items(), key=lambda kv: kv[1], reverse=True)[:max_items]:
        parts.append(f"U+{ord(ch):04X}:{cnt}")
    more = len(stats) - min(len(stats), max_items)
    if more > 0:
        parts.append(f"+{more} weitere")
    return f"Repl: {total} (" + ", ".join(parts) + ")"

def _format_stats_multiline(stats: Counter[str]) -> str:
    """Mehrzeilig, ausführlich."""
    total = sum(stats.values())
    if total <= 0:
        return "  (keine Ersetzungen)\n"
    out = []
    for ch, cnt in sorted(stats.items(), key=lambda kv: kv[1], reverse=True):
        name = unicodedata.name(ch, "<unbekannt>")
        repl = REPLACEMENTS.get(ch, "")
        out.append(f"  U+{ord(ch):04X} {name}: {cnt}  →  {repr(repl)}")
    out.append(f"  Summe Ersetzungen: {total}")
    return "\n".join(out) + "\n"

# ---------- .docx ----------
def clean_docx(input_file: Path, output_file: Path, stats: Optional[Counter[str]] = None) -> Tuple[bool, str]:
    if not HAVE_PYTHON_DOCX:
        return False, "python-docx nicht installiert (pip install python-docx)."
    try:
        doc = Document(str(input_file))
        changed = False

        for para in doc.paragraphs:
            for run in para.runs:
                new = clean_text(run.text, stats)
                if new != run.text:
                    run.text = new
                    changed = True

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        for run in para.runs:
                            new = clean_text(run.text, stats)
                            if new != run.text:
                                run.text = new
                                changed = True

        output_file.parent.mkdir(parents=True, exist_ok=True)
        if changed:
            doc.save(str(output_file))
            return True, "ok"
        else:
            if output_file != input_file and not output_file.exists():
                shutil.copy2(str(input_file), str(output_file))
            return True, "unverändert"
    except Exception as e:
        return False, f"Fehler: {e}"

# ---------- Heuristik Textdatei ----------
def probably_text(path: Path, max_bytes: int = 8192) -> bool:
    try:
        with open(path, "rb") as f:
            chunk = f.read(max_bytes)
        if b"\x00" in chunk:
            return False
        non_print = sum(b < 9 or (13 < b < 32) for b in chunk)
        return non_print / max(1, len(chunk)) < 0.3
    except Exception:
        return False

TEXT_EXTENSIONS_COMMON = {
    ".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".log", ".xml", ".html", ".htm",
    ".tex", ".bib"
}

CANDIDATE_ENCODINGS = ("utf-8", "utf-8-sig", "cp1252", "latin-1")

def read_text_with_guess(path: Path) -> Tuple[Optional[str], Optional[str]]:
    data = path.read_bytes()
    for enc in CANDIDATE_ENCODINGS:
        try:
            return data.decode(enc), enc
        except UnicodeDecodeError:
            continue
    try:
        return data.decode("utf-8", errors="replace"), "utf-8"
    except Exception:
        return None, None

def detect_newline_style(text: str) -> str | None:
    if "\r\n" in text:
        return "\r\n"
    if "\r" in text and "\n" not in text:
        return "\r"
    if "\n" in text:
        return "\n"
    return None

def write_text_preserving_newlines(path: Path, text: str, newline: str | None) -> None:
    with open(path, "w", encoding="utf-8", newline=newline) as f:
        f.write(text)

def should_treat_as_text(path: Path, force_all_text: bool) -> bool:
    if path.suffix.lower() in TEXT_EXTENSIONS_COMMON:
        return True
    if path.suffix.lower() == ".py":
        return True
    if force_all_text:
        return probably_text(path)
    return False

def clean_textfile(
    input_file: Path,
    output_file: Path,
    validate_python: bool = False,
    stats: Optional[Counter[str]] = None
) -> Tuple[bool, str]:
    try:
        original, _enc = read_text_with_guess(input_file)
        if original is None:
            return False, "Konnte Text nicht lesen/decodieren."
        newline = detect_newline_style(original)
        cleaned = clean_text(original, stats)
        changed = cleaned != original

        if validate_python and input_file.suffix.lower() == ".py":
            import ast
            try:
                ast.parse(original)
            except SyntaxError as e:
                return False, f"Ursprüngliche Python-Datei fehlerhaft: {e}"
            try:
                ast.parse(cleaned)
            except SyntaxError as e:
                return False, f"Bereinigung abgebrochen: Syntax würde brechen ({e})."

        output_file.parent.mkdir(parents=True, exist_ok=True)
        if changed:
            write_text_preserving_newlines(output_file, cleaned, newline)
            return True, "ok"
        else:
            if output_file != input_file and not output_file.exists():
                shutil.copy2(str(input_file), str(output_file))
            return True, "unverändert"
    except Exception as e:
        return False, f"Fehler: {e}"

# ---------- .pptx ----------
def _clean_textframe(tf, stats: Optional[Counter[str]] = None) -> bool:
    changed = False
    for p in tf.paragraphs:
        for r in p.runs:
            new = clean_text(r.text, stats)
            if new != r.text:
                r.text = new
                changed = True
    return changed

def clean_pptx(input_file: Path, output_file: Path, stats: Optional[Counter[str]] = None) -> Tuple[bool, str]:
    if not HAVE_PYTHON_PPTX:
        return False, "python-pptx nicht installiert (pip install python-pptx)."
    try:
        prs = Presentation(str(input_file))
        changed = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if _clean_textframe(shape.text_frame, stats):
                        changed = True
                if hasattr(shape, "has_table") and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame:
                                if _clean_textframe(cell.text_frame, stats):
                                    changed = True
            if slide.has_notes_slide:
                notes = slide.notes_slide
                if hasattr(notes, "notes_text_frame") and notes.notes_text_frame:
                    if _clean_textframe(notes.notes_text_frame, stats):
                        changed = True
                else:
                    for shape in notes.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            if _clean_textframe(shape.text_frame, stats):
                                changed = True

        output_file.parent.mkdir(parents=True, exist_ok=True)
        if changed:
            prs.save(str(output_file))
            return True, "ok"
        else:
            if output_file != input_file and not output_file.exists():
                shutil.copy2(str(input_file), str(output_file))
            return True, "unverändert"
    except Exception as e:
        return False, f"Fehler: {e}"

# ---------- .xlsx ----------
def clean_xlsx(input_file: Path, output_file: Path, stats: Optional[Counter[str]] = None) -> Tuple[bool, str]:
    if not HAVE_OPENPYXL:
        return False, "openpyxl nicht installiert (pip install openpyxl)."
    try:
        wb = openpyxl.load_workbook(str(input_file))
        changed = False

        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    v = cell.value
                    if isinstance(v, str) and not v.startswith("="):
                        new = clean_text(v, stats)
                        if new != v:
                            cell.value = new
                            changed = True
                    if cell.comment and cell.comment.text:
                        newc = clean_text(cell.comment.text, stats)
                        if newc != cell.comment.text:
                            from openpyxl.comments import Comment
                            cell.comment = Comment(newc, cell.comment.author)
                            changed = True

        output_file.parent.mkdir(parents=True, exist_ok=True)
        if changed:
            wb.save(str(output_file))
            return True, "ok"
        else:
            if output_file != input_file and not output_file.exists():
                shutil.copy2(str(input_file), str(output_file))
            return True, "unverändert"
    except Exception as e:
        return False, f"Fehler: {e}"

# ---------- .pdf ----------
def clean_pdf(input_file: Path, output_file: Path, stats: Optional[Counter[str]] = None) -> Tuple[bool, str]:
    """
    PDF-Verarbeitung:
    - extrahiert Text (Copy/Paste Text) aus allen Seiten
    - bereinigt unsichtbare Unicode-Zeichen
    - schreibt Ergebnis als *_cleaned.txt ins Zielverzeichnis

    Hinweis: Das PDF selbst wird NICHT überschrieben (Layout- und Inhaltsstreams sind komplex).
    """
    if not HAVE_PYMUPDF:
        return False, "PyMuPDF nicht installiert (pip install pymupdf)."
    try:
        doc = fitz.open(str(input_file))
        parts: list[str] = []
        for page in doc:
            parts.append(page.get_text("text"))

        original = "\n".join(parts)
        cleaned = clean_text(original, stats)

        # Ziel: immer Textdatei erzeugen
        out_txt = output_file
        if out_txt.suffix.lower() == ".pdf":
            out_txt = out_txt.with_name(out_txt.stem + "_cleaned.txt")
        elif out_txt.suffix.lower() != ".txt":
            out_txt = out_txt.with_suffix(out_txt.suffix + "_cleaned.txt")

        out_txt.parent.mkdir(parents=True, exist_ok=True)
        out_txt.write_text(cleaned, encoding="utf-8")

        if cleaned != original:
            return True, f"ok (Text exportiert: {out_txt.name})"
        return True, f"unverändert (Text exportiert: {out_txt.name})"
    except Exception as e:
        return False, f"Fehler: {e}"

# ---------- Dispatcher ----------
def process_one(
    path: Path,
    out_dir: Optional[Path],
    in_place: bool,
    backup: bool,
    validate_python: bool,
    force_all_text: bool,
    stats: Optional[Counter[str]] = None
) -> Tuple[bool, str]:
    # Automatisch clean_out neben der Datei anlegen, falls kein Ausgabeverzeichnis definiert
    if out_dir is None and not in_place:
        out_dir = path.parent / "clean_out"
        out_dir.mkdir(parents=True, exist_ok=True)

    target = path if in_place else (out_dir / path.name)

    if backup and in_place:
        bak = path.with_suffix(path.suffix + ".bak")
        try:
            if not bak.exists():
                shutil.copy2(str(path), str(bak))
        except Exception:
            pass

    ext = path.suffix.lower()

    if ext == ".docx":
        return clean_docx(path, target, stats)
    if ext == ".pptx":
        return clean_pptx(path, target, stats)
    if ext == ".xlsx":
        return clean_xlsx(path, target, stats)
    if ext == ".pdf":
        # PDF wird als bereinigter Text exportiert
        return clean_pdf(path, target, stats)

    if should_treat_as_text(path, force_all_text):
        return clean_textfile(path, target, validate_python=validate_python, stats=stats)

    return False, "Übersprungen (nicht als unterstütztes Format / Text erkannt)."

def iter_files(inputs: Iterable[Path], recursive: bool) -> Iterable[Path]:
    for p in inputs:
        if p.is_dir():
            if recursive:
                for f in p.rglob("*"):
                    if f.is_file():
                        yield f
            else:
                for f in p.glob("*"):
                    if f.is_file():
                        yield f
        elif p.is_file():
            yield p

# ---------- CLI ----------
def parse_args(argv: list[str]) -> argparse.Namespace:
    ap = argparse.ArgumentParser(
        description="Entfernt unsichtbare Unicode-Steuerzeichen aus DOCX, PPTX, XLSX, PDF (Text-Export) sowie Text- und Python-Dateien. Mit GUI."
    )
    ap.add_argument("paths", nargs="*", help="Dateien oder Ordner")
    group = ap.add_mutually_exclusive_group()
    group.add_argument("-i", "--in-place", action="store_true", help="Dateien an Ort und Stelle überschreiben")
    group.add_argument("-o", "--out-dir", type=Path, help="Ausgabeverzeichnis (Struktur wird beibehalten)")
    ap.add_argument("-r", "--recursive", action="store_true", help="Ordner rekursiv verarbeiten")
    ap.add_argument("--backup", action="store_true", help="Vor Überschreiben .bak-Backup anlegen (nur mit --in-place)")
    ap.add_argument("--validate-python", action="store_true", help="Vor/Nach-Bereinigung Python-Syntax prüfen (.py)")
    ap.add_argument("--force-all-text", action="store_true", help="Auch unbekannte Endungen als Text versuchen (Heuristik)")
    ap.add_argument("--gui", action="store_true", help="GUI starten")
    return ap.parse_args(argv)

# ---------- GUI (Tkinter) ----------
def launch_gui():
    import tkinter as tk
    from tkinter import ttk, filedialog, messagebox

    class App(tk.Tk):
        def __init__(self):
            super().__init__()
            self.title("Unicode Cleaner")
            self.geometry("820x560")
            self.minsize(740, 520)

            # State
            self.inputs: List[Path] = []
            self.out_dir: Optional[Path] = None
            self.in_place = tk.BooleanVar(value=False)
            self.recursive = tk.BooleanVar(value=True)
            self.backup = tk.BooleanVar(value=True)
            self.validate_python = tk.BooleanVar(value=False)
            self.force_all_text = tk.BooleanVar(value=False)
            self.running = False
            self.cancel_requested = False

            self._build_ui()

            # Worker comms
            self.q = queue.Queue()

            # Dependencies hint
            missing = []
            if not HAVE_PYTHON_DOCX:  missing.append("python-docx")
            if not HAVE_PYTHON_PPTX:  missing.append("python-pptx")
            if not HAVE_OPENPYXL:     missing.append("openpyxl")
            if not HAVE_PYMUPDF:      missing.append("pymupdf")
            if missing:
                self._append_log(
                    "Hinweis: Für einige Formate bitte installieren:\n  pip install " + " ".join(missing) + "\n\n"
                )

        def _build_ui(self):
            pad = {"padx": 10, "pady": 8}

            frm_top = ttk.Frame(self)
            frm_top.pack(fill="x", **pad)

            ttk.Label(frm_top, text="Ausgewählte Pfade:").grid(row=0, column=0, sticky="w")

            self.lst = tk.Listbox(frm_top, height=6)
            self.lst.grid(row=1, column=0, columnspan=4, sticky="nsew", pady=(4, 6))
            frm_top.columnconfigure(0, weight=1)

            btn_add_files = ttk.Button(frm_top, text="Dateien hinzufügen", command=self._add_files)
            btn_add_dir   = ttk.Button(frm_top, text="Ordner hinzufügen", command=self._add_dir)
            btn_clear     = ttk.Button(frm_top, text="Liste leeren", command=self._clear_list)
            btn_add_files.grid(row=2, column=0, sticky="w")
            btn_add_dir.grid(row=2, column=1, sticky="w")
            btn_clear.grid(row=2, column=2, sticky="w")

            frm_opts = ttk.LabelFrame(self, text="Optionen")
            frm_opts.pack(fill="x", **pad)

            ttk.Checkbutton(frm_opts, text="In-place überschreiben", variable=self.in_place, command=self._toggle_outdir).grid(row=0, column=0, sticky="w")
            ttk.Checkbutton(frm_opts, text="Rekursiv (bei Ordnern)", variable=self.recursive).grid(row=0, column=1, sticky="w")
            ttk.Checkbutton(frm_opts, text="Backup (.bak) anlegen", variable=self.backup).grid(row=0, column=2, sticky="w")
            ttk.Checkbutton(frm_opts, text="Python-Syntax prüfen (.py)", variable=self.validate_python).grid(row=1, column=0, sticky="w")
            ttk.Checkbutton(frm_opts, text="Unbekannte Endungen als Text versuchen", variable=self.force_all_text).grid(row=1, column=1, sticky="w")

            self.outdir_lbl = ttk.Label(frm_opts, text="Ausgabeverzeichnis:")
            self.outdir_lbl.grid(row=2, column=0, sticky="w", pady=(8, 0))
            self.outdir_val = tk.StringVar(value=str(Path.cwd() / "cleaned_output"))
            self.out_dir = Path(self.outdir_val.get())
            self.ent_outdir = ttk.Entry(frm_opts, textvariable=self.outdir_val, width=62)
            self.ent_outdir.grid(row=2, column=1, sticky="w", pady=(8,0))
            ttk.Button(frm_opts, text="Wählen…", command=self._choose_outdir).grid(row=2, column=2, sticky="w", pady=(8,0))

            for i in range(3):
                frm_opts.columnconfigure(i, weight=1)

            frm_run = ttk.Frame(self)
            frm_run.pack(fill="x", **pad)

            self.btn_start = ttk.Button(frm_run, text="Start", command=self._start)
            self.btn_cancel = ttk.Button(frm_run, text="Abbrechen", command=self._cancel, state="disabled")
            self.btn_start.pack(side="left")
            self.btn_cancel.pack(side="left", padx=8)

            self.prog = ttk.Progressbar(frm_run, mode="determinate")
            self.prog.pack(fill="x", padx=10, expand=True)

            frm_log = ttk.LabelFrame(self, text="Log")
            frm_log.pack(fill="both", expand=True, **pad)

            self.txt = tk.Text(frm_log, wrap="none", height=14)
            self.txt.pack(fill="both", expand=True)
            self.txt.configure(state="disabled")

        def _toggle_outdir(self):
            disabled = self.in_place.get()
            state = "disabled" if disabled else "normal"
            self.ent_outdir.configure(state=state)
            if disabled:
                self._append_log("Hinweis: In-place aktiv – Ausgabeverzeichnis wird ignoriert.\n")

        def _add_files(self):
            paths = filedialog.askopenfilenames(
                title="Dateien auswählen",
                filetypes=[
                    ("Unterstützt", "*.docx *.pdf *.pptx *.xlsx *.txt *.md *.csv *.tsv *.json *.xml *.html *.py"),
                    ("Word (.docx)", "*.docx"),
                    ("PDF (.pdf)", "*.pdf"),
                    ("Alle Dateien", "*.*"),
                ]
            )
            for p in paths:
                self.inputs.append(Path(p))
                self.lst.insert("end", p)

        def _add_dir(self):
            p = filedialog.askdirectory(title="Ordner auswählen")
            if p:
                self.inputs.append(Path(p))
                self.lst.insert("end", p)

        def _clear_list(self):
            self.inputs.clear()
            self.lst.delete(0, "end")

        def _choose_outdir(self):
            p = filedialog.askdirectory(title="Ausgabeverzeichnis wählen")
            if p:
                self.outdir_val.set(p)
                self.out_dir = Path(p)

        def _append_log(self, s: str):
            self.txt.configure(state="normal")
            self.txt.insert("end", s)
            self.txt.see("end")
            self.txt.configure(state="disabled")

        def _set_running(self, running: bool):
            self.running = running
            self.btn_start.configure(state="disabled" if running else "normal")
            self.btn_cancel.configure(state="normal" if running else "disabled")

        def _cancel(self):
            if self.running:
                self.cancel_requested = True
                self._append_log("Abbruch angefordert…\n")

        def _start(self):
            if not self.inputs:
                messagebox.showwarning("Hinweis", "Bitte mindestens eine Datei / einen Ordner hinzufügen.")
                return

            # Resolve out_dir
            if not self.in_place.get():
                self.out_dir = Path(self.outdir_val.get()).resolve()
                try:
                    self.out_dir.mkdir(parents=True, exist_ok=True)
                except Exception as e:
                    messagebox.showerror("Fehler", f"Ausgabeverzeichnis kann nicht erstellt werden:\n{e}")
                    return
            else:
                self.out_dir = None

            # Worker vorbereiten
            self._set_running(True)
            self.cancel_requested = False
            self.q = queue.Queue()
            self.prog["value"] = 0

            # Dateien zählen (für Progress)
            files = list(iter_files(self.inputs, self.recursive.get()))
            total = len(files)
            if total == 0:
                messagebox.showinfo("Info", "Keine verarbeitbaren Dateien gefunden.")
                self._set_running(False)
                return
            self.prog["maximum"] = total

            # Thread starten
            t = threading.Thread(
                target=self._worker,
                args=(files, self.out_dir, self.in_place.get(), self.backup.get(), self.validate_python.get(), self.force_all_text.get()),
                daemon=True
            )
            t.start()
            self.after(100, self._poll_queue)

        def _worker(self, files: List[Path], out_dir: Optional[Path], in_place: bool, backup: bool, validate_python: bool, force_all_text: bool):
            processed = changed = skipped = failed = 0
            idx = 0
            stats_total: Counter[str] = Counter()

            for f in files:
                if self.cancel_requested:
                    self.q.put(("log", "\n>> Abgebrochen.\n"))
                    break

                stats_file: Counter[str] = Counter()
                ok, msg = process_one(
                    f, out_dir, in_place, backup, validate_python, force_all_text,
                    stats=stats_file
                )
                if stats_file:
                    stats_total.update(stats_file)

                idx += 1
                if ok:
                    processed += 1
                    if msg == "ok" or msg.startswith("ok "):
                        changed += 1
                else:
                    if msg.startswith("Übersprungen"):
                        skipped += 1
                    else:
                        failed += 1

                self.q.put(("step", idx, f, ok, msg, stats_file))

            self.q.put(("done", processed, changed, skipped, failed, stats_total))

        def _poll_queue(self):
            try:
                while True:
                    item = self.q.get_nowait()
                    if item[0] == "step":
                        _, idx, f, ok, msg, stats_file = item
                        self.prog["value"] = idx
                        statline = _format_stats_one_liner(stats_file)
                        self._append_log(f"[{'OK' if ok else '!!'}] {f}: {msg} | {statline}\n")

                    elif item[0] == "log":
                        _, s = item
                        self._append_log(s)

                    elif item[0] == "done":
                        _, processed, changed, skipped, failed, stats_total = item

                        self._append_log("\nZusammenfassung:\n")
                        self._append_log(f"  Verarbeitet: {processed}\n  Geändert:    {changed}\n  Übersprungen:{skipped}\n  Fehler:      {failed}\n")

                        self._append_log("\nErsetzungsstatistik (gesamt):\n")
                        self._append_log(_format_stats_multiline(stats_total))

                        self._set_running(False)
                        if self.cancel_requested:
                            messagebox.showinfo("Abgebrochen", "Vorgang wurde abgebrochen.")
                        else:
                            if failed == 0:
                                messagebox.showinfo("Fertig", "Bereinigung abgeschlossen.")
                            else:
                                messagebox.showwarning("Fertig mit Fehlern", f"Bereinigung abgeschlossen – Fehler: {failed}")
                    self.q.task_done()
            except queue.Empty:
                pass
            if self.running:
                self.after(100, self._poll_queue)

    # DPI-Awareness für scharfe Fonts auf Windows
    try:
        import ctypes
        ctypes.windll.shcore.SetProcessDpiAwareness(1)
    except Exception:
        pass

    app = App()
    app.mainloop()

# ---------- main ----------
def main(argv: list[str]) -> int:
    args = parse_args(argv)

    # GUI starten, wenn:
    # 1) --gui angegeben ist, oder
    # 2) keine Pfade via CLI übergeben wurden (typischer Doppelklick)
    if args.gui or not args.paths:
        launch_gui()
        return 0

    paths = [Path(p).resolve() for p in args.paths]

    if (not args.in_place) and (args.out_dir is None):
        print("Hinweis: Ohne --in-place wird in ein Ausgabeverzeichnis geschrieben.")
        args.out_dir = Path.cwd() / "cleaned_output"
        args.out_dir.mkdir(parents=True, exist_ok=True)

    processed = changed = skipped = failed = 0
    stats_total: Counter[str] = Counter()

    files = list(iter_files(paths, args.recursive))
    total = len(files)
    if total == 0:
        print("Keine verarbeitbaren Dateien gefunden.")
        return 0

    for i, f in enumerate(files, 1):
        stats_file: Counter[str] = Counter()
        ok, msg = process_one(
            f,
            args.out_dir,
            args.in_place,
            args.backup,
            args.validate_python,
            args.force_all_text,
            stats=stats_file
        )
        if stats_file:
            stats_total.update(stats_file)

        if ok:
            processed += 1
            if msg == "ok" or msg.startswith("ok "):
                changed += 1
        else:
            if msg.startswith("Übersprungen"):
                skipped += 1
            else:
                failed += 1

        print(f"[{i:>4}/{total}] [{'OK' if ok else '!!'}] {f}: {msg} | {_format_stats_one_liner(stats_file)}")

    print("\nZusammenfassung:")
    print(f"  Verarbeitet: {processed}")
    print(f"  Geändert:    {changed}")
    print(f"  Übersprungen:{skipped}")
    print(f"  Fehler:      {failed}")

    print("\nErsetzungsstatistik (gesamt):")
    sys.stdout.write(_format_stats_multiline(stats_total))

    return 0 if failed == 0 else 2

if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
